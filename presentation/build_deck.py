"""Builds the ScreenSmart pitch deck (team ctrl alt elite) for the Garaza FinTech AI
Hackathon. Run: pip install python-pptx && python build_deck.py
(run render_svg.py first to produce architecture.png + garaza_logo.png)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os, random

# ---- palette (matches the ScreenSmart dashboard) ----
BG     = RGBColor.from_string("0B0F1A")
PANEL  = RGBColor.from_string("141C2B")
PANEL2 = RGBColor.from_string("0F1626")
BORDER = RGBColor.from_string("2A3550")
TEAL   = RGBColor.from_string("5EEAD4")
BLUE   = RGBColor.from_string("38BDF8")
VIOLET = RGBColor.from_string("818CF8")
GREEN  = RGBColor.from_string("28E68C")
AMBER  = RGBColor.from_string("F5C32D")
RED    = RGBColor.from_string("F74856")
TEXT   = RGBColor.from_string("E8EDF6")
MUTED  = RGBColor.from_string("8B95A9")
FAINT  = RGBColor.from_string("5B647A")
WHITE  = RGBColor.from_string("FFFFFF")
FONT   = "Segoe UI"

# dim "node" colours for the background dot-field (echoes the live dashboard)
DIM = [RGBColor.from_string(h) for h in
       ("173A2A", "16403D", "153049", "3E3318", "3E1C24", "262247", "1d2740", "1d2740")]
POP = [GREEN, AMBER, RED, TEAL, BLUE]

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
FRONT = os.path.join(ROOT, "docs", "front.png")
ARCH = os.path.join(HERE, "architecture.png")
LOGO = os.path.join(HERE, "garaza_logo.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
rng = random.Random(11)


def _noshadow(o):
    o.shadow.inherit = False
    return o


def box(s, shp, l, t, w, h, fill=PANEL, line=BORDER, line_w=1.0):
    o = s.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    o.fill.solid(); o.fill.fore_color.rgb = fill
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = line; o.line.width = Pt(line_w)
    return _noshadow(o)


def text(s, l, t, w, h, runs, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, color, bold, size)]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    for txt, c, b, sz in runs:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = font
    return tb


def dotfield(s, n=44):
    """Scatter a dim node-field across the slide (drawn behind content)."""
    for _ in range(n):
        x = rng.uniform(0.15, 13.0); y = rng.uniform(0.25, 7.2)
        if rng.random() < 0.18:
            d = rng.uniform(0.08, 0.15); c = rng.choice(POP)
        else:
            d = rng.uniform(0.05, 0.13); c = rng.choice(DIM)
        box(s, MSO_SHAPE.OVAL, x, y, d, d, fill=c, line=None)


def slide(dots=True):
    s = prs.slides.add_slide(BLANK)
    box(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=BG, line=None)
    if dots:
        dotfield(s)
    return s


def kicker(s, emoji, txt, color=TEAL):
    text(s, 0.7, 0.5, 9, 0.4, [(emoji + "  ", color, True, 14), (txt.upper(), color, True, 13)])


def title(s, txt):
    text(s, 0.7, 0.92, 12, 1.0, txt, size=34, color=WHITE, bold=True)


def footer(s):
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(11.45), Inches(6.95), height=Inches(0.3))
    text(s, 0.7, 6.96, 7, 0.35, [("ScreenSmart", TEXT, True, 10), ("  ·  ctrl alt elite", FAINT, False, 10)])


# ============================================================ 1 · TITLE
s = slide()
text(s, 0, 2.15, 13.333, 1.4, [("Screen", WHITE, True, 80), ("Smart", TEAL, True, 80)], align=PP_ALIGN.CENTER)
text(s, 0, 3.78, 13.333, 0.6, [("🛡️  ", TEAL, False, 18), ("Real-time sanctions screening — every payment, sub-second.", MUTED, False, 20)], align=PP_ALIGN.CENTER)
box(s, MSO_SHAPE.RECTANGLE, 5.67, 4.55, 2.0, 0.03, fill=TEAL, line=None)
text(s, 0, 4.78, 13.333, 0.5, [("👥  TEAM   ", FAINT, True, 16), ("ctrl alt elite", WHITE, True, 19)], align=PP_ALIGN.CENTER)
# Garaza logo + event line at the bottom
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches((13.333 - 1.9) / 2), Inches(6.35), width=Inches(1.9))
text(s, 0, 6.92, 13.333, 0.4, "FinTech AI Hackathon · Belgrade · 2026", size=12, color=FAINT, align=PP_ALIGN.CENTER)

# ============================================================ 2 · AGENDA
s = slide()
kicker(s, "🧭", "Contents", BLUE)
title(s, "What’s inside")
agenda = [
    ("01", "🎯", "The problem", "why sanctions are easy to dodge", BLUE),
    ("02", "📊", "Fun facts", "financial crime, by the numbers", AMBER),
    ("03", "🛠️", "ScreenSmart", "what we built & the tech stack", TEAL),
    ("04", "🕸️", "Architecture", "streaming scatter-gather over Kafka", VIOLET),
    ("05", "🚀", "Live demo", "the dashboard in action", GREEN),
]
y = 2.25
for num, emo, head, sub, col in agenda:
    box(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y, 11.5, 0.82, fill=PANEL, line=BORDER)
    text(s, 1.15, y + 0.1, 1.0, 0.6, num, size=24, color=col, bold=True)
    text(s, 2.15, y, 0.8, 0.82, emo, size=22, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.0, y, 4.6, 0.82, head, size=19, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.6, y, 5.6, 0.82, sub, size=15, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.95
footer(s)

# ============================================================ 3 · ABOUT US
s = slide()
kicker(s, "🙌", "Who we are", TEAL)
title(s, "Team — ctrl alt elite")
for cx, initials, name, col in [(3.0, "MS", "Mateja Subin", BLUE), (8.0, "AM", "Andrija Milikić", VIOLET)]:
    box(s, MSO_SHAPE.OVAL, cx - 1.1, 2.35, 2.2, 2.2, fill=PANEL2, line=col, line_w=2.5)
    text(s, cx - 1.1, 2.35, 2.2, 2.2, initials, size=40, color=col, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx - 2.0, 4.8, 4.0, 0.5, name, size=22, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    text(s, cx - 2.0, 5.35, 4.0, 0.4, [("💻  builder · ctrl alt elite", MUTED, False, 14)], align=PP_ALIGN.CENTER)
text(s, 0, 6.3, 13.333, 0.5, "Two builders at the intersection of AI and finance.", size=14, color=FAINT, align=PP_ALIGN.CENTER)
footer(s)

# ============================================================ 4 · PROBLEM
s = slide()
kicker(s, "🎯", "The problem", RED)
title(s, "Sanctions are easy to dodge")
text(s, 0.7, 1.8, 11.8, 0.8,
     "A blacklist check on a name catches the obvious. Real evasion hides one step away — "
     "behind shell companies, nominee owners and chains of intermediaries.", size=17, color=MUTED)
probs = [
    ("🏢", "Shell & front companies", "a sanctioned party routes money through a clean-looking company it secretly controls", BLUE),
    ("🔁", "Layering & hops", "funds pass through several accounts so no single name matches a list", VIOLET),
    ("🔤", "Typos & transliteration", "“Qadhafi” vs “Kaddafi” — fuzzy spelling slips past exact matching", TEAL),
    ("🪙", "Crypto & cross-border", "wallets and corridors a static list never sees", AMBER),
]
y = 2.9
for emo, head, body, col in probs:
    box(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y, 11.5, 0.78, fill=PANEL, line=BORDER)
    box(s, MSO_SHAPE.OVAL, 1.1, y + 0.16, 0.46, 0.46, fill=PANEL2, line=col, line_w=1.5)
    text(s, 1.1, y + 0.13, 0.46, 0.5, emo, size=16, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.8, y, 3.6, 0.78, head, size=16, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.4, y, 6.8, 0.78, body, size=13.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.86
box(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y + 0.05, 11.5, 0.6, fill=PANEL2, line=RED, line_w=1.2)
text(s, 1.15, y + 0.05, 11.0, 0.6, [("⚠️  The gap:  ", RED, True, 14), ("a name match isn’t enough — you have to follow the network.", TEXT, False, 14)], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ============================================================ 5 · FUN FACTS
s = slide()
kicker(s, "📊", "Fun facts", AMBER)
title(s, "Financial crime, by the numbers")
cards = [
    ("🌍", "2–5%", "of global GDP is laundered every year — up to $2 trillion", AMBER),
    ("🕵️", "~1%", "of that laundered money is ever detected or seized", RED),
    ("🏢", "$70B", "estimated yearly cost of shell-company abuse to the US alone", BLUE),
    ("📂", "29,000", "shell companies exposed in the Pandora Papers leak", VIOLET),
    ("📈", "100+", "names added to OFAC’s sanctions list 8 separate times in 2024", TEAL),
    ("💸", "$3T+", "in illicit funds flowed through the financial system in 2023", GREEN),
]
cw, ch, gx, gy = 3.75, 1.5, 0.3, 0.35
x0, y0 = 0.9, 2.2
for i, (emo, big, sub, col) in enumerate(cards):
    cx = x0 + (i % 3) * (cw + gx); cy = y0 + (i // 3) * (ch + gy)
    box(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch, fill=PANEL, line=BORDER)
    box(s, MSO_SHAPE.RECTANGLE, cx, cy, 0.08, ch, fill=col, line=None)
    text(s, cx + 0.22, cy + 0.1, 0.7, 0.6, emo, size=22)
    text(s, cx + 0.95, cy + 0.1, cw - 1.1, 0.62, big, size=29, color=col, bold=True)
    text(s, cx + 0.27, cy + 0.78, cw - 0.45, 0.65, sub, size=12, color=MUTED, spacing=1.05)
text(s, 0.9, 6.95, 11.8, 0.35, "Sources: UNODC · OCCRP · FinCEN · ICIJ Pandora Papers · CNAS 2024 · AML Intelligence", size=9.5, color=FAINT)

# ============================================================ 6 · TECH STACK
s = slide()
kicker(s, "🛠️", "How it’s built", TEAL)
title(s, "Tech stack")
rows = [
    ("Screening engine", "🧠", BLUE,   ["🐍 Python", "🌳 LightGBM", "📈 scikit-learn", "🔤 RapidFuzz"]),
    ("Streaming",        "⚡", VIOLET, ["🟪 Apache Kafka", "🔁 aiokafka", "🧩 scatter-gather"]),
    ("Data & API",       "🗄️", AMBER,  ["🐘 PostgreSQL", "🚀 FastAPI", "🔐 JWT", "🔌 WebSocket"]),
    ("Frontend",         "🖥️", GREEN,  ["⚛️ React", "⚡ Vite", "🎨 live canvas"]),
    ("Infra",            "🐳", TEAL,   ["🐳 Docker", "🌐 nginx", "📦 Compose"]),
]
y = 2.2
for name, emo, col, chips in rows:
    box(s, MSO_SHAPE.OVAL, 0.9, y + 0.05, 0.6, 0.6, fill=PANEL2, line=col, line_w=1.8)
    text(s, 0.9, y + 0.02, 0.6, 0.62, emo, size=20, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.65, y, 2.5, 0.7, name, size=15.5, color=col, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    cx = 4.25
    for chip in chips:
        w = 0.42 + 0.108 * len(chip)
        box(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, y + 0.12, w, 0.46, fill=PANEL, line=BORDER)
        text(s, cx, y + 0.12, w, 0.46, chip, size=12.5, color=TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w + 0.18
    y += 0.9
footer(s)

# ============================================================ 7 · ARCHITECTURE
s = slide()
kicker(s, "🕸️", "Architecture", VIOLET)
title(s, "Scatter-gather over Kafka")
if os.path.exists(ARCH):
    img_w = Inches(11.6)
    s.shapes.add_picture(ARCH, int((SW - img_w) / 2), Inches(1.95), width=img_w)
text(s, 0, 6.95, 13.333, 0.4,
     "Two engines screen every payment in parallel — name/identity + network exposure — and the accumulator escalates to the worst verdict.",
     size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ============================================================ 8 · DEMO
s = slide()
kicker(s, "🚀", "Let’s screen some payments", GREEN)
title(s, "Demo time")
if os.path.exists(FRONT):
    img_w = Inches(10.6)
    s.shapes.add_picture(FRONT, int((SW - img_w) / 2), Inches(2.0), width=img_w)
text(s, 0, 6.7, 13.333, 0.5, [("🟢 live feed   ", GREEN, True, 13), ("· click any dot · analyst review desk · network-exposure drill-down", MUTED, False, 13)], align=PP_ALIGN.CENTER)
footer(s)

out = os.path.join(HERE, "ScreenSmart_CtrlAltElite.pptx")
prs.save(out)
print("saved", out)
